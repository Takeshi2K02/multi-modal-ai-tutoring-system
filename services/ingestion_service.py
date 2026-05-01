import os
import uuid
import json
from datetime import datetime
from typing import List, Dict, Any
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from services.vector_factory import get_vector_db

RAW_DOC_DIR = os.path.join(os.getcwd(), "local_data", "raw_documents")
PROCESSED_CHUNKS_DIR = os.path.join(os.getcwd(), "local_data", "processed_chunks")

# Ensure dirs exist
os.makedirs(RAW_DOC_DIR, exist_ok=True)
os.makedirs(PROCESSED_CHUNKS_DIR, exist_ok=True)

from bson import ObjectId
from db.connection import get_db_connection

async def ingest_document(file_content: bytes, filename: str, collection_id: str = None, plan_id: str = None):
    """
    Process any supported document file and ingest into VectorDB.
    """
    if not collection_id:
        collection_id = str(uuid.uuid4())
        print(f"[Ingest] Generated new collection_id (UUID): {collection_id}")

    print(f"Ingesting {filename} (Collection: {collection_id})...")
    
    # 1. Save Raw File
    file_path = os.path.join(RAW_DOC_DIR, filename)
    with open(file_path, "wb") as f:
        f.write(file_content)
        
    # 2. Extract Text & Chunk based on extension
    ext = os.path.splitext(filename)[1].lower()
    
    if ext == ".pdf":
        chunks = _process_pdf_text(file_path, filename)
    elif ext == ".docx":
        chunks = _process_docx_text(file_path, filename)
    elif ext == ".pptx":
        chunks = _process_pptx_text(file_path, filename)
    elif ext == ".txt":
        chunks = _process_txt_text(file_path, filename)
    else:
        raise ValueError(f"Unsupported file format: {ext}")
    
    # 3. Save Parsed Chunks (Debug)
    debug_path = os.path.join(PROCESSED_CHUNKS_DIR, f"{filename}.json")
    with open(debug_path, "w") as f:
        json.dump(chunks, f, indent=2)
    
    # 4. Push to VectorDB
    vectordb = get_vector_db()
    
    docs_to_add = []
    for chunk in chunks:
        meta = {
            "source_file": filename,
            "lecture_title": chunk["lecture_title"],
            "page_number": chunk.get("page_number", 1),
            "chunk_index": chunk["chunk_index"],
            "file_type": ext
        }
        if collection_id:
            meta["collection_id"] = collection_id
            print(f"[Ingest] Assigned collection_id: {collection_id} to {filename}")
        else:
            print(f"[Ingest] WARNING: No collection_id provided for {filename}")

        docs_to_add.append({
            "id": chunk["doc_id"],
            "text": chunk["text"],
            "metadata": meta
        })
        
    if docs_to_add:
        print(f"[Ingest] Sample Metadata (Chunk 0): {docs_to_add[0]['metadata']}")
    vectordb.add_documents(docs_to_add)

    # --- BUG 1: Write collection_id back to Plan Document ---
    if plan_id:
        try:
            db = get_db_connection()
            db.learning_plans.update_one(
                { "_id": ObjectId(plan_id) },
                { "$set": { "system_metadata.collection_id": collection_id } }
            )
            print(f"[Ingest] ✅ collection_id {collection_id} saved to plan {plan_id}")
        except Exception as e:
            print(f"[Ingest] ❌ Failed to save collection_id to plan: {e}")

    return {"status": "success", "chunks_count": len(chunks), "collection_id": collection_id}

def _chunk_text(text: str, doc_id_prefix: str, lecture_title: str, filename: str, page_num: int = 1) -> List[Dict[str, Any]]:
    """Helper to chunk text consistently."""
    CHUNK_SIZE = 500
    OVERLAP = 50
    words = text.split()
    chunks = []
    chunk_counter = 0
    
    for j in range(0, len(words), CHUNK_SIZE - OVERLAP):
        chunk_words = words[j:j + CHUNK_SIZE]
        chunk_text = " ".join(chunk_words)
        
        if len(chunk_text) < 5:
            continue
            
        chunks.append({
            "doc_id": f"{doc_id_prefix}_c{chunk_counter}",
            "source_file": filename,
            "lecture_title": lecture_title,
            "page_number": page_num,
            "chunk_index": chunk_counter,
            "text": chunk_text
        })
        chunk_counter += 1
    return chunks

def _process_pdf_text(file_path: str, filename: str) -> List[Dict[str, Any]]:
    reader = PdfReader(file_path)
    all_chunks = []
    lecture_title = filename.replace(".pdf", "").replace("_", " ").title()
    
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if not text:
            continue
        all_chunks.extend(_chunk_text(text, f"{filename}_p{i}", lecture_title, filename, i + 1))
            
    return all_chunks

def _process_docx_text(file_path: str, filename: str) -> List[Dict[str, Any]]:
    doc = Document(file_path)
    text = "\n".join([p.text for p in doc.paragraphs])
    lecture_title = filename.replace(".docx", "").replace("_", " ").title()
    return _chunk_text(text, filename, lecture_title, filename)

def _process_pptx_text(file_path: str, filename: str) -> List[Dict[str, Any]]:
    prs = Presentation(file_path)
    all_chunks = []
    lecture_title = filename.replace(".pptx", "").replace("_", " ").title()
    
    for i, slide in enumerate(prs.slides):
        text_runs = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
        text = "\n".join(text_runs)
        if text.strip():
            all_chunks.extend(_chunk_text(text, f"{filename}_s{i}", lecture_title, filename, i + 1))
            
    return all_chunks

def _process_txt_text(file_path: str, filename: str) -> List[Dict[str, Any]]:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    lecture_title = filename.replace(".txt", "").replace("_", " ").title()
    return _chunk_text(text, filename, lecture_title, filename)
